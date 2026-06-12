#!/usr/bin/env python3
"""Generate an EDITABLE, banded PowerPoint deck from a module's `session:` data.

Native python-pptx shapes (bands, panels, pills, callouts) in BioFAIR colours,
so the .pptx looks like the colleague's deck *and* stays fully editable.
One source (the module front-matter) → this deck, the .docx guide, the web deck.

Run:  scripts/.venv/bin/python scripts/session_to_pptx.py
"""
import os, re, sys, glob, subprocess, tempfile
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.abspath('.')
MODULES = os.path.join(ROOT, 'src/content/modules')
PUBLIC = os.path.join(ROOT, 'public')

# ELIXIR-UK palette (names kept for stability; values are ELIXIR brand)
INK     = RGBColor(0x02, 0x34, 0x52)  # navy — structure / band / text
GREEN   = RGBColor(0xF4, 0x7D, 0x20)  # orange — accent / chip / rule
GREEN_D = RGBColor(0xB5, 0x61, 0x10)  # dark orange
MAGENTA = RGBColor(0x03, 0x7E, 0xAB)  # mid blue — secondary accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x4A, 0x65, 0x77)  # slate
PANEL_L = RGBColor(0xDD, 0xE3, 0xE7)  # rule
LGREEN  = RGBColor(0xFD, 0xF0, 0xE4)  # light orange tint
LMAGENTA= RGBColor(0xE6, 0xF2, 0xF7)  # light blue tint
SUBTLE  = RGBColor(0xC7, 0xD2, 0xDA)
FONT = 'Calibri'  # Word/PowerPoint-safe until the ELIXIR (Lato) reference template is supplied

W, H = Inches(13.333), Inches(7.5)


def clean(s=''):
    s = re.sub(r'<code>(.*?)</code>', r'\1', s or '')
    s = re.sub(r'<[^>]+>', '', s)
    return s.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>').strip()


def duration(minutes):
    # `minutes` is a cumulative range like "40–50"; show how long the SECTION
    # takes (the duration), which is what a facilitator actually needs.
    nums = re.findall(r'\d+', minutes or '')
    if len(nums) >= 2:
        return f"{int(nums[1]) - int(nums[0])} min"
    return f"{minutes} min" if minutes else ''


def frontmatter(path):
    raw = open(path, encoding='utf-8').read()
    m = re.match(r'^---\n(.*?)\n---', raw, re.S)
    return yaml.safe_load(m.group(1)) if m else None


def rasterize(src):
    p = os.path.join(PUBLIC, src.lstrip('/'))
    if not os.path.exists(p):
        return None
    if not p.endswith('.svg'):
        return p
    out = os.path.join(tempfile.mkdtemp(), 'img.png')
    subprocess.run(['rsvg-convert', '-w', '1400', p, '-o', out], check=True)
    return out


def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    if radius is not None:
        try: sp.adjustments[0] = radius
        except Exception: pass
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, runs, size, color, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=FONT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = runs if isinstance(runs, list) else [runs]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = font; r.font.color.rgb = color
    return tb


def pill(slide, x, y, label):
    w = Inches(max(0.55, 0.135 * len(label) + 0.32)); h = Inches(0.34)
    sp = rect(slide, x, y, w, h, RGBColor(0xF1, 0xF3, 0xF2), line=INK,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(11); r.font.bold = True
    r.font.name = FONT; r.font.color.rgb = INK
    return Emu(x + w)


def bullets(slide, x, y, w, h, points):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run(); r.text = '•  ' + clean(pt)
        r.font.size = Pt(18); r.font.name = FONT; r.font.color.rgb = INK


def build(prs, fm):
    blank = prs.slide_layouts[6]

    # ── Title slide ──
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = INK
    logo = os.path.join(PUBLIC, 'assets/branding/BioFAIR-logo.png')
    if os.path.exists(logo):
        chip = rect(s, Inches(0.7), Inches(0.65), Inches(2.7), Inches(0.95), WHITE)
        s.shapes.add_picture(logo, Inches(0.85), Inches(0.78), height=Inches(0.7))
    total = (fm['session'][-1].get('minutes') or '').split('–')[-1]
    text(s, Inches(0.72), Inches(2.0), Inches(11), Inches(0.4),
         f"WORKSHOP · MODULE 02 · {len(fm['session'])} STEPS · ~{total} MIN", 13, GREEN, bold=True)
    text(s, Inches(0.7), Inches(2.45), Inches(11.5), Inches(1.8),
         clean(fm['title']), 40, WHITE, bold=True)
    rect(s, Inches(0.75), Inches(4.5), Inches(1.0), Inches(0.09), GREEN)
    rect(s, Inches(1.75), Inches(4.5), Inches(1.0), Inches(0.09), MAGENTA)
    text(s, Inches(0.72), Inches(4.8), Inches(10.5), Inches(1.6),
         clean(fm['summary']), 16, SUBTLE)
    text(s, Inches(0.72), Inches(6.85), Inches(11.5), Inches(0.4),
         'Built by the ELIXIR-UK RDM Club  ·  commissioned by BioFAIR', 12, GREEN, bold=True)

    # ── One slide per step ──
    for i, step in enumerate(fm['session']):
        s = prs.slides.add_slide(blank)
        # band
        rect(s, 0, 0, W, Inches(1.05), INK)
        text(s, Inches(0.5), Inches(0.1), Inches(9.6), Inches(0.85),
             clean(step['title']), 26, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        if step.get('minutes'):
            chip = rect(s, Inches(11.0), Inches(0.32), Inches(1.95), Inches(0.42), GREEN,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
            tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = duration(step['minutes']); r.font.size = Pt(12)
            r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE
        has_img = bool(step.get('image'))
        has_ex = bool(step.get('exercise'))
        panel_w = Inches(5.6) if has_img else Inches(12.33)
        panel_h = Inches(3.6) if has_ex else Inches(5.0)
        panel_y = Inches(1.35)
        # panel + green accent
        rect(s, Inches(0.5), panel_y, panel_w, panel_h, WHITE, line=PANEL_L)
        rect(s, Inches(0.5), panel_y, panel_w, Inches(0.09), GREEN)
        bullets(s, Inches(0.55), Emu(panel_y + Inches(0.18)),
                Emu(panel_w - Inches(0.1)), Emu(panel_h - Inches(0.3)), step.get('points') or [])
        # image (two-column)
        if has_img:
            img = rasterize(step['image']['src'])
            if img:
                s.shapes.add_picture(img, Inches(6.4), Inches(2.0), width=Inches(6.5))
        # exercise callout
        if has_ex:
            ey = Inches(5.25); eh = Inches(1.35)
            rect(s, Inches(0.5), ey, Inches(0.08), eh, MAGENTA)
            rect(s, Inches(0.58), ey, Inches(12.25), eh, LMAGENTA)
            text(s, Inches(0.8), ey, Inches(11.85), eh,
                 '✏  Exercise: ' + clean(step['exercise']), 16, INK, anchor=MSO_ANCHOR.MIDDLE)
        # footer + logo
        text(s, Inches(0.5), Inches(7.08), Inches(10), Inches(0.3),
             'Built by the ELIXIR-UK RDM Club  ·  Module 02 — FAIRification of RNA-seq', 11, MUTED)
        if os.path.exists(logo):
            s.shapes.add_picture(logo, Inches(12.0), Inches(6.95), height=Inches(0.4))
        # speaker notes = the facilitator script
        if step.get('script'):
            s.notes_slide.notes_text_frame.text = clean(step['script'])


def main():
    n = 0
    for path in sorted(glob.glob(os.path.join(MODULES, '*.md'))):
        if os.path.basename(path).startswith('_'):
            continue
        fm = frontmatter(path)
        if not fm or not fm.get('session'):
            continue
        mod_id = os.path.basename(path)[:-3]
        sub = mod_id if mod_id.startswith('p2-') else mod_id[:2]
        out_dir = os.path.join(PUBLIC, 'assets/modules', sub)
        os.makedirs(out_dir, exist_ok=True)
        prs = Presentation(); prs.slide_width = W; prs.slide_height = H
        build(prs, fm)
        out = os.path.join(out_dir, f'{mod_id}.pptx')
        prs.save(out)
        print(f'✓ {mod_id}: {os.path.relpath(out, ROOT)} ({len(prs.slides._sldIdLst)} slides)')
        n += 1
    print(f'Generated {n} banded deck(s).' if n else 'No modules with a session: block.')


if __name__ == '__main__':
    main()
